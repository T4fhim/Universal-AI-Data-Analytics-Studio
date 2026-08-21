# File: src/readers/powerpoint_reader.py
"""Reads tables embedded in PowerPoint (.pptx) slides into a Dataset.

Multi-table like :class:`~src.readers.excel_reader.ExcelReader`, but
the "tables" here are PowerPoint table shapes scattered across slides
rather than worksheets — a deck can have zero, one, or many table
shapes, in any slide, so :meth:`list_tables` walks every slide once to
enumerate them rather than assuming a fixed structure the way a
workbook's sheet list is fixed. Each table's first row is assumed to
be its header, matching this project's existing header-row convention
for :class:`~src.readers.excel_reader.ExcelReader` — the same
documented limitation applies (a title row above the real header is
not detected).
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.exc import PythonPptxError

from src.core.exceptions import ReaderError
from src.core.logger import get_logger
from src.readers.base_reader import BaseReader
from src.services.workspace_service import Dataset

_logger = get_logger(__name__)

_POWERPOINT_EXTENSIONS = {".pptx"}


class PowerPointReader(BaseReader):
    """Reads table shapes found on any slide of a .pptx presentation."""

    SUPPORTED_EXTENSIONS = _POWERPOINT_EXTENSIONS

    @classmethod
    def can_read(cls, path: Path) -> bool:
        return path.suffix.lower() in _POWERPOINT_EXTENSIONS

    @classmethod
    def list_tables(cls, path: Path) -> list[str]:
        """Return one label per table shape found across every slide, e.g. ``"Slide 2, Table 1"``.

        Raises:
            ReaderError: If the file does not exist, cannot be opened
                as a .pptx presentation, or contains no table shapes at
                all.
        """
        table_shapes = cls._find_table_shapes(path)
        if not table_shapes:
            raise ReaderError(f"No tables found in any slide of {path}.")
        return [label for label, _shape in table_shapes]

    @classmethod
    def read(cls, path: Path, table_name: str | None = None) -> Dataset:
        """Read one table shape from the presentation at ``path``.

        Args:
            path: The presentation to read.
            table_name: Which table to read, by the label returned
                from :meth:`list_tables` (e.g. ``"Slide 2, Table 1"``).
                If the presentation has exactly one table shape,
                ``None`` reads it directly — same single-table
                convenience as
                :meth:`~src.readers.excel_reader.ExcelReader.read`.

        Raises:
            ReaderError: If the file does not exist, cannot be opened,
                has more than one table but no ``table_name`` was
                given, ``table_name`` does not match any table found,
                or the selected table has no rows.
        """
        table_shapes = cls._find_table_shapes(path)
        if not table_shapes:
            raise ReaderError(f"No tables found in any slide of {path}.")

        labels = [label for label, _shape in table_shapes]
        if table_name is None:
            if len(table_shapes) == 1:
                table_name = labels[0]
            else:
                raise ReaderError(
                    f"{path} contains {len(table_shapes)} table(s) "
                    f"({', '.join(labels)}); specify which one to read "
                    f"via the table_name argument."
                )
        elif table_name not in labels:
            raise ReaderError(
                f"{path} has no table labeled '{table_name}'. "
                f"Available tables: {', '.join(labels)}."
            )

        shape = dict(table_shapes)[table_name]
        rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
        if len(rows) < 2:
            raise ReaderError(
                f"Table '{table_name}' in {path} has no data rows beyond its header."
            )

        dataframe = pd.DataFrame(rows[1:], columns=rows[0])

        _logger.info(
            "Read PowerPoint table '%s' from %s: %d rows, %d columns.",
            table_name,
            path,
            len(dataframe),
            len(dataframe.columns),
        )

        return Dataset(
            name=f"{path.stem} — {table_name}",
            dataframe=dataframe,
            source_format="pptx",
            source_path=path,
        )

    @classmethod
    def _find_table_shapes(cls, path: Path) -> list[tuple[str, object]]:
        """Return ``(label, shape)`` pairs for every table shape found across every slide, in order."""
        if not path.exists():
            raise ReaderError(f"PowerPoint file does not exist: {path}")

        try:
            presentation = Presentation(str(path))
        except (PythonPptxError, OSError, KeyError) as exc:
            raise ReaderError(
                f"Failed to open {path} as a PowerPoint file: {exc}"
            ) from exc

        found: list[tuple[str, object]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            table_index_on_slide = 0
            for shape in slide.shapes:
                if shape.has_table:
                    table_index_on_slide += 1
                    found.append(
                        (f"Slide {slide_index}, Table {table_index_on_slide}", shape)
                    )
        return found
