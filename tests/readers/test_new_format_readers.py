# File: tests/readers/test_new_format_readers.py
"""Tests for the seven milestone-14 readers: ODS, YAML, Parquet, Feather, PowerPoint, HTML, ZIP/GZIP.

Registry-level concerns (dispatch by extension, no extension
collisions, plugin registration) are already covered generically for
every registered reader — including these seven — by
tests/readers/test_reader_registry.py's ``@pytest.mark.parametrize``
over ``_BUILTIN_READERS``; this module instead covers each new
reader's own format-specific parsing and error paths, which that
generic test cannot exercise.
"""

from __future__ import annotations

import gzip
import zipfile
from pathlib import Path

import pandas as pd
import pytest

from src.core.exceptions import ReaderError


@pytest.fixture()
def sample_dataframe() -> pd.DataFrame:
    return pd.DataFrame({"a": [1, 2], "b": [3, 4]})


# -- OdsReader ----------------------------------------------------------


def test_ods_reader_reads_a_single_sheet(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.ods_reader import OdsReader

    path = tmp_path / "data.ods"
    sample_dataframe.to_excel(path, engine="odf", index=False)

    dataset = OdsReader.read(path)

    assert dataset.row_count == 2
    assert dataset.column_count == 2
    assert dataset.source_format == "ods"


def test_ods_reader_missing_file_raises(tmp_path: Path) -> None:
    from src.readers.ods_reader import OdsReader

    with pytest.raises(ReaderError):
        OdsReader.read(tmp_path / "does_not_exist.ods")


# -- YamlReader -----------------------------------------------------------


def test_yaml_reader_reads_a_list_of_records(tmp_path: Path) -> None:
    from src.readers.yaml_reader import YamlReader

    path = tmp_path / "data.yaml"
    path.write_text("- a: 1\n  b: 3\n- a: 2\n  b: 4\n", encoding="utf-8")

    dataset = YamlReader.read(path)

    assert dataset.row_count == 2
    assert list(dataset.dataframe.columns) == ["a", "b"]
    assert dataset.source_format == "yaml"


def test_yaml_reader_rejects_a_scalar_top_level_document(tmp_path: Path) -> None:
    from src.readers.yaml_reader import YamlReader

    path = tmp_path / "data.yaml"
    path.write_text("42\n", encoding="utf-8")

    with pytest.raises(ReaderError):
        YamlReader.read(path)


def test_yaml_reader_empty_file_raises(tmp_path: Path) -> None:
    from src.readers.yaml_reader import YamlReader

    path = tmp_path / "empty.yaml"
    path.write_text("", encoding="utf-8")

    with pytest.raises(ReaderError):
        YamlReader.read(path)


# -- ParquetReader ----------------------------------------------------------


def test_parquet_reader_reads_a_file(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.parquet_reader import ParquetReader

    path = tmp_path / "data.parquet"
    sample_dataframe.to_parquet(path)

    dataset = ParquetReader.read(path)

    assert dataset.row_count == 2
    assert dataset.source_format == "parquet"


def test_parquet_reader_invalid_file_raises(tmp_path: Path) -> None:
    from src.readers.parquet_reader import ParquetReader

    path = tmp_path / "not_parquet.parquet"
    path.write_text("this is not parquet", encoding="utf-8")

    with pytest.raises(ReaderError):
        ParquetReader.read(path)


# -- FeatherReader ----------------------------------------------------------


def test_feather_reader_reads_a_file(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.feather_reader import FeatherReader

    path = tmp_path / "data.feather"
    sample_dataframe.to_feather(path)

    dataset = FeatherReader.read(path)

    assert dataset.row_count == 2
    assert dataset.source_format == "feather"


def test_feather_reader_invalid_file_raises(tmp_path: Path) -> None:
    from src.readers.feather_reader import FeatherReader

    path = tmp_path / "not_feather.feather"
    path.write_text("this is not feather", encoding="utf-8")

    with pytest.raises(ReaderError):
        FeatherReader.read(path)


# -- HtmlReader ----------------------------------------------------------


def test_html_reader_reads_a_single_table(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.html_reader import HtmlReader

    path = tmp_path / "page.html"
    path.write_text(sample_dataframe.to_html(index=False), encoding="utf-8")

    dataset = HtmlReader.read(path)

    assert dataset.row_count == 2
    assert dataset.source_format == "html"


def test_html_reader_multiple_tables_requires_table_name(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.html_reader import HtmlReader

    path = tmp_path / "page.html"
    html = sample_dataframe.to_html(index=False) + sample_dataframe.to_html(index=False)
    path.write_text(html, encoding="utf-8")

    tables = HtmlReader.list_tables(path)
    assert tables == ["Table 1", "Table 2"]

    with pytest.raises(ReaderError):
        HtmlReader.read(path)

    dataset = HtmlReader.read(path, table_name="Table 2")
    assert dataset.row_count == 2


def test_html_reader_no_tables_raises(tmp_path: Path) -> None:
    from src.readers.html_reader import HtmlReader

    path = tmp_path / "page.html"
    path.write_text(
        "<html><body><p>No tables here.</p></body></html>", encoding="utf-8"
    )

    with pytest.raises(ReaderError):
        HtmlReader.read(path)


# -- PowerPointReader ----------------------------------------------------------


def _make_pptx_with_table(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(
        3, 2, Inches(1), Inches(1), Inches(4), Inches(2)
    )
    table = table_shape.table
    table.cell(0, 0).text = "a"
    table.cell(0, 1).text = "b"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "3"
    table.cell(2, 0).text = "2"
    table.cell(2, 1).text = "4"
    presentation.save(str(path))


def test_powerpoint_reader_reads_a_single_table(tmp_path: Path) -> None:
    from src.readers.powerpoint_reader import PowerPointReader

    path = tmp_path / "deck.pptx"
    _make_pptx_with_table(path)

    dataset = PowerPointReader.read(path)

    assert dataset.row_count == 2
    assert list(dataset.dataframe.columns) == ["a", "b"]
    assert dataset.source_format == "pptx"


def test_powerpoint_reader_no_tables_raises(tmp_path: Path) -> None:
    from pptx import Presentation

    from src.readers.powerpoint_reader import PowerPointReader

    path = tmp_path / "empty_deck.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(str(path))

    with pytest.raises(ReaderError):
        PowerPointReader.read(path)


# -- ArchiveReader (ZIP/GZIP) ----------------------------------------------------------


def test_archive_reader_reads_a_single_file_zip(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.archive_reader import ArchiveReader

    inner_csv = tmp_path / "inner.csv"
    sample_dataframe.to_csv(inner_csv, index=False)
    zip_path = tmp_path / "archive.zip"
    with zipfile.ZipFile(zip_path, "w") as archive:
        archive.write(inner_csv, arcname="inner.csv")

    dataset = ArchiveReader.read(zip_path)

    assert dataset.row_count == 2
    assert dataset.name == "inner"


def test_archive_reader_zip_with_multiple_readable_files_requires_table_name(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.archive_reader import ArchiveReader

    inner_csv = tmp_path / "inner.csv"
    sample_dataframe.to_csv(inner_csv, index=False)
    inner_json = tmp_path / "inner.json"
    inner_json.write_text(sample_dataframe.to_json(orient="records"), encoding="utf-8")

    zip_path = tmp_path / "archive.zip"
    with zipfile.ZipFile(zip_path, "w") as archive:
        archive.write(inner_csv, arcname="inner.csv")
        archive.write(inner_json, arcname="inner.json")

    names = ArchiveReader.list_tables(zip_path)
    assert set(names) == {"inner.csv", "inner.json"}

    with pytest.raises(ReaderError):
        ArchiveReader.read(zip_path)

    dataset = ArchiveReader.read(zip_path, table_name="inner.csv")
    assert dataset.row_count == 2


def test_archive_reader_zip_excludes_unreadable_inner_files(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.archive_reader import ArchiveReader

    inner_csv = tmp_path / "inner.csv"
    sample_dataframe.to_csv(inner_csv, index=False)
    inner_exe = tmp_path / "not_data.exe"
    inner_exe.write_bytes(b"\x00\x01\x02")

    zip_path = tmp_path / "archive.zip"
    with zipfile.ZipFile(zip_path, "w") as archive:
        archive.write(inner_csv, arcname="inner.csv")
        archive.write(inner_exe, arcname="not_data.exe")

    names = ArchiveReader.list_tables(zip_path)
    assert names == ["inner.csv"]


def test_archive_reader_reads_a_gzip_file(
    tmp_path: Path, sample_dataframe: pd.DataFrame
) -> None:
    from src.readers.archive_reader import ArchiveReader

    inner_csv = tmp_path / "inner.csv"
    sample_dataframe.to_csv(inner_csv, index=False)
    gz_path = tmp_path / "inner.csv.gz"
    with open(inner_csv, "rb") as source, gzip.open(gz_path, "wb") as destination:
        destination.write(source.read())

    dataset = ArchiveReader.read(gz_path)

    assert dataset.row_count == 2
    assert dataset.name == "inner"


def test_archive_reader_bad_zip_raises(tmp_path: Path) -> None:
    from src.readers.archive_reader import ArchiveReader

    path = tmp_path / "bad.zip"
    path.write_bytes(b"not a real zip file")

    with pytest.raises(ReaderError):
        ArchiveReader.read(path)


def test_archive_reader_missing_file_raises(tmp_path: Path) -> None:
    from src.readers.archive_reader import ArchiveReader

    with pytest.raises(ReaderError):
        ArchiveReader.read(tmp_path / "does_not_exist.zip")
